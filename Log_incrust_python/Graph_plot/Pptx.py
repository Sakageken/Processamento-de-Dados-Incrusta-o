import pandas as pd
import numpy as np
import pptx
import io
from dateutil.parser import parse
from pathlib import Path
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import matplotlib.pyplot as plt
from scipy.fft import fft # Using scipy's fft
from Graph_plot import graphs_plot, Data_processing, search
from PIL import Image
from datetime import datetime
from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas


def pptx_generation(figures, df_table, file_to_read, path_txt_folder):

        pptx_name = parse(str(file_to_read), dayfirst=True, fuzzy=True)
        input_format_1 = "%Y-%m-%d"
        output_format_1 = "%d_%m_%Y"
        date_obj_1 = datetime.strptime(str(pptx_name.date()), input_format_1)
        pptx_name = date_obj_1.strftime(output_format_1)
        


        path_presentation = Path(path_txt_folder / "pptx_files")
        if path_presentation.is_dir():
             pass
        else:
             path_presentation.mkdir(parents=True, exist_ok=True)

        presentation = pptx.Presentation()
        slide_layout = presentation.slide_layouts[6]
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        #Saving the figures temporally
        figure_ids = plt.get_fignums()

        print(figure_ids)
        print(f"Found {len(figure_ids)} figures to process: {figure_ids}")
        
        for fig_id in range(1,len(figure_ids)):
            # --- Activate the figure ---
            fig = plt.figure(fig_id)


        print(figure_ids)
        print(f"Found {len(figure_ids)} figures to process: {figure_ids}")
            
        for fig_id in range(1,len(figure_ids)):
            # --- Activate the figure ---
            fig = plt.figure(fig_id)


            # --- Save the active figure to an in-memory buffer ---
            image_stream = io.BytesIO()
            fig.savefig(image_stream, format='png', dpi=300)
            image_stream.seek(0)

            # --- Add a new slide and place the image on it ---
            slide = presentation.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = f'Plot from Figure {fig_id}'

            width = slide_width
            height = slide_height
            top = (slide_height - height) // 2
            left = (slide_width - width) // 2
            slide.shapes.add_picture(image_stream,left,top,width=width,height=height)
            image_stream.close()

        #Slide 6
        slide_6 = presentation.slides.add_slide(slide_layout)
        rows, cols = df_table.shape[0] + 1, df_table.shape[1] + 1 # +1 for header row and +1 for the row label
        width = (slide_width//10) * 8
        height = (slide_height//10) * 8
        top = (slide_height - height) // 2
        left = (slide_width - width) // 2

        table_shape = slide_6.shapes.add_table(rows,cols,left,top,width,height)
        table = table_shape.table
        column_fraction = [3/6,1/6,2/6]
        for column, fraction in zip(table.columns, column_fraction):
            column.width = int((width * fraction))

        table.cell(0,0).text = df_table.index.name
        table.cell(0,0).vertical_anchor = MSO_ANCHOR.MIDDLE
        table.cell(0,0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for col_index, col_name in enumerate(df_table.columns):
            # The column index in the pptx table is offset by 1
            table.cell(0, col_index + 1).text = col_name
            table.cell(0, col_index + 1).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(0, col_index + 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 5. Populate the Data Rows
        for row_index, (index_label, row_data) in enumerate(df_table.iterrows()):
            # Populate the first column with the index label
            # The row index in the pptx table is offset by 1 for the header
            table.cell(row_index + 1, 0).text = str(index_label)
            table.cell(row_index + 1, 0).text_frame.word_wrap = False
            table.cell(row_index + 1, 0).vertical_anchor = MSO_ANCHOR.MIDDLE
            table.cell(row_index + 1, 0).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Populate the rest of the columns with the row's data
            for col_index, cell_value in enumerate(row_data):
                # Both row and column indexes are offset by 1
                table.cell(row_index + 1, col_index + 1).text = str(cell_value)
                table.cell(row_index + 1, col_index + 1).vertical_anchor = MSO_ANCHOR.MIDDLE
                table.cell(row_index + 1, col_index + 1).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        presentation_name = f"Gráficos teste {pptx_name}.pptx"
        #presentation_name = f"Gráficos teste {new_date_string_1}.pptx"
        presentation_path = path_presentation / presentation_name
        presentation.save(presentation_path)
        print("Presentation with DataFrame table created successfully! 📊")

if __name__ == "__main__":
    pptx_generation
    print("File should not be run as main")
